"""Geometry QA for the deck, since this machine has no LibreOffice to render with.

Catches the defects a visual pass would: anything past the slide edge, images stretched
out of their true aspect ratio, and text boxes overlapping a shape they were meant to sit
beside rather than on.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

W, H = 13.333, 7.5
MARGIN = 0.45

def inches(v): return Emu(v).inches if v is not None else None

prs = Presentation(sys.argv[1] if len(sys.argv) > 1 else "presentation/ChartQA-Week1.pptx")
problems = []
for n, slide in enumerate(prs.slides, 1):
    boxes = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        x, y = inches(sh.left), inches(sh.top)
        w, h = inches(sh.width), inches(sh.height)
        name = (sh.text_frame.text.strip().split("\n")[0][:38]
                if sh.has_text_frame and sh.text_frame.text.strip() else sh.shape_type)
        if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
            problems.append(f"slide {n}: OFF-SLIDE  ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}  {name}")
        elif x < MARGIN or y < MARGIN or x + w > W - MARGIN or y + h > H - MARGIN:
            edge = min(x, y, W - (x + w), H - (y + h))
            if edge < MARGIN - 0.12:
                problems.append(f"slide {n}: tight margin {edge:.2f}in  {name}")
        boxes.append((x, y, w, h, name, sh.shape_type))
    # images must keep their true aspect ratio
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.image is not None:
            iw, ih = sh.image.size
            want, got = iw / ih, inches(sh.width) / inches(sh.height)
            if abs(want - got) / want > 0.02:
                problems.append(f"slide {n}: IMAGE STRETCHED {got:.3f} vs true {want:.3f}")

    # A text block that *partially* overlaps a filled card is the defect a render would
    # show: fully inside is the intended case, fully outside is fine, half-in is a bug.
    cards = [b for b in boxes if str(b[5]).startswith("AUTO_SHAPE")]
    texts = [b for b in boxes if str(b[5]).startswith("TEXT_BOX")]
    for cx, cy, cw, ch, cname, _ in cards:
        for tx, ty, tw, th, tname, _ in texts:
            ox = min(cx + cw, tx + tw) - max(cx, tx)
            oy = min(cy + ch, ty + th) - max(cy, ty)
            if ox <= 0.02 or oy <= 0.02:
                continue
            inside = (tx >= cx - 0.02 and ty >= cy - 0.02
                      and tx + tw <= cx + cw + 0.02 and ty + th <= cy + ch + 0.02)
            if not inside:
                spill = max(cy - ty, ty + th - (cy + ch), cx - tx, tx + tw - (cx + cw))
                if spill > 0.06:
                    problems.append(
                        f"slide {n}: text spills {spill:.2f}in past its card  "
                        f"'{tname}' vs '{cname}'")

print(f"{len(prs.slides)} slides checked")
for p in problems:
    print(" ", p)
print("  no geometry problems" if not problems else f"  {len(problems)} to fix")
